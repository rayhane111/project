from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from transformers import pipeline
import textwrap
import fitz  # PyMuPDF pour PDF
from docx import Document
import openpyxl  # Pour Excel
from pptx import Presentation
from functools import lru_cache

app = FastAPI()

# Servir les fichiers statiques (comme index.html)
app.mount("/static", StaticFiles(directory="static"), name="static")

@app.get("/", response_class=HTMLResponse)
async def read_root():
    with open("static/index.html", "r", encoding="utf-8") as file:
        return HTMLResponse(content=file.read())

# Dictionnaire des codes de langues (ISO 639-1)
LANGUAGE_CODES = {
    "Anglais": "en",
    "Francais": "fr",
    "Arabe": "ar",
    "Espagnol": "es",
}

# Vérification de la disponibilité des modèles
AVAILABLE_MODELS = {
    "fr-en": "Helsinki-NLP/opus-mt-fr-en",
    "en-fr": "Helsinki-NLP/opus-mt-en-fr",
    "ar-en": "Helsinki-NLP/opus-mt-ar-en",
    "en-ar": "Helsinki-NLP/opus-mt-en-ar",
    "es-en": "Helsinki-NLP/opus-mt-es-en",
    "en-es": "Helsinki-NLP/opus-mt-en-es",
}

# Cache pour éviter de recharger les modèles à chaque requête
@lru_cache(maxsize=10)
def load_translator(src_code: str, tgt_code: str):
    model_key = f"{src_code}-{tgt_code}"
    
    if model_key in AVAILABLE_MODELS:
        model_name = AVAILABLE_MODELS[model_key]
        print(f"🔹 Chargement du modèle : {model_name}")
        return pipeline("translation", model=model_name)

    # Gestion des traductions indirectes via l'anglais
    elif src_code == "fr" and tgt_code == "ar":
        print("🔹 Traduction indirecte via l'anglais : fr -> en -> ar")
        return (
            pipeline("translation", model=AVAILABLE_MODELS["fr-en"]),
            pipeline("translation", model=AVAILABLE_MODELS["en-ar"])
        )

    elif src_code == "ar" and tgt_code == "fr":
        print("🔹 Traduction indirecte via l'anglais : ar -> en -> fr")
        return (
            pipeline("translation", model=AVAILABLE_MODELS["ar-en"]),
            pipeline("translation", model=AVAILABLE_MODELS["en-fr"])
        )

    else:
        raise ValueError(f"⚠️ Aucun modèle disponible pour {src_code} -> {tgt_code}")

# Découpe le texte en segments de taille gérable
def chunk_text(text, max_length=400):
    return textwrap.wrap(text, max_length)

# Extraction du texte en fonction du type de fichier
def extract_text(file: UploadFile):
    try:
        if file.filename.endswith(".txt"):
            return file.file.read().decode("utf-8")

        elif file.filename.endswith(".pdf"):
            doc = fitz.open(stream=file.file.read(), filetype="pdf")
            return "\n".join([page.get_text() for page in doc])

        elif file.filename.endswith(".docx"):
            doc = Document(file.file)
            return "\n".join([para.text for para in doc.paragraphs])

        elif file.filename.endswith(".xlsx"):
            wb = openpyxl.load_workbook(file.file)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows():
                    text += "\t".join([str(cell.value or "") for cell in row]) + "\n"
            return text

        elif file.filename.endswith(".pptx"):
            prs = Presentation(file.file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            raise HTTPException(status_code=400, detail="❌ Type de fichier non supporté.")

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erreur lors de l'extraction du texte : {str(e)}")

@app.post("/upload/")
async def upload_file(
    file: UploadFile = File(...), 
    src_lang: str = Form(...), 
    tgt_lang: str = Form(...)
):
    text = extract_text(file)

    if not text.strip():
        raise HTTPException(status_code=400, detail="❌ Aucun texte extrait du fichier.")

    src_code = LANGUAGE_CODES.get(src_lang)
    tgt_code = LANGUAGE_CODES.get(tgt_lang)

    if not src_code or not tgt_code:
        raise HTTPException(status_code=400, detail=f"⚠️ Langue non supportée : {src_lang} -> {tgt_lang}")

    try:
        # Charger le modèle correct
        translator = load_translator(src_code, tgt_code)

        # Si traduction indirecte via l'anglais
        if isinstance(translator, tuple):
            translator1, translator2 = translator
            intermediate_text = "\n".join([translator1(chunk)[0]['translation_text'] for chunk in chunk_text(text)])
            translated_text = "\n".join([translator2(chunk)[0]['translation_text'] for chunk in chunk_text(intermediate_text)])
        else:
            translated_text = "\n".join([translator(chunk)[0]['translation_text'] for chunk in chunk_text(text)])

        return {"translated_text": translated_text}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"❌ Erreur interne : {str(e)}")
