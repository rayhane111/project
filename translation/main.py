from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from transformers import pipeline
import textwrap
import fitz  # PyMuPDF for PDF handling
from docx import Document
import openpyxl  # For Excel
from pptx import Presentation
from functools import lru_cache
import os

# Initialize FastAPI app
app = FastAPI()

# Set the correct path for static files
STATIC_DIR = r"C:\Users\User\doc_translation_service\translation\static"

# Ensure the static directory exists
if not os.path.exists(STATIC_DIR):
    os.makedirs(STATIC_DIR)

# Mount static files (serves index.html)
app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")

@app.get("/", response_class=HTMLResponse)
async def read_root():
    index_path = os.path.join(STATIC_DIR, "index.html")
    try:
        with open(index_path, "r", encoding="utf-8") as file:
            return HTMLResponse(content=file.read())
    except FileNotFoundError:
        raise HTTPException(status_code=404, detail="index.html not found in static folder.")

# Language codes (ISO 639-1)
LANGUAGE_CODES = {
    "Anglais": "en",
    "Francais": "fr",
    "Arabe": "ar",
    "Espagnol": "es",
}

# Available translation models
AVAILABLE_MODELS = {
    "fr-en": "Helsinki-NLP/opus-mt-fr-en",
    "en-fr": "Helsinki-NLP/opus-mt-en-fr",
    "ar-en": "Helsinki-NLP/opus-mt-ar-en",
    "en-ar": "Helsinki-NLP/opus-mt-en-ar",
    "es-en": "Helsinki-NLP/opus-mt-es-en",
    "en-es": "Helsinki-NLP/opus-mt-en-es",
}

# Cache model loading
@lru_cache(maxsize=10)
def load_translator(src_code: str, tgt_code: str):
    model_key = f"{src_code}-{tgt_code}"
    
    if model_key in AVAILABLE_MODELS:
        return pipeline("translation", model=AVAILABLE_MODELS[model_key])

    elif src_code != "en" and tgt_code != "en":
        return (
            pipeline("translation", model=AVAILABLE_MODELS.get(f"{src_code}-en")),
            pipeline("translation", model=AVAILABLE_MODELS.get(f"en-{tgt_code}"))
        )

    else:
        raise ValueError(f"No model available for {src_code} -> {tgt_code}")

# Split text into chunks
def chunk_text(text, max_length=400):
    return textwrap.wrap(text, max_length)

# Extract text based on file type
def extract_text(file: UploadFile):
    try:
        if file.filename.endswith(".txt"):
            return file.file.read().decode("utf-8")

        elif file.filename.endswith(".pdf"):
            doc = fitz.open(stream=file.file.read(), filetype="pdf")
            return "\n".join([page.get_text() for page in doc])

        elif file.filename.endswith(".docx"):
            doc = Document(file.file)
            return "\n".join([para.text for para in doc.paragraphs])

        elif file.filename.endswith(".xlsx"):
            wb = openpyxl.load_workbook(file.file)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows():
                    text += "\t".join([str(cell.value or "") for cell in row]) + "\n"
            return text

        elif file.filename.endswith(".pptx"):
            prs = Presentation(file.file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            raise HTTPException(status_code=400, detail="File type not supported.")

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text: {str(e)}")

@app.post("/upload/")
async def upload_file(
    file: UploadFile = File(...), 
    src_lang: str = Form(...), 
    tgt_lang: str = Form(...)
):
    text = extract_text(file)

    if not text.strip():
        raise HTTPException(status_code=400, detail="No text extracted from the file.")

    src_code = LANGUAGE_CODES.get(src_lang)
    tgt_code = LANGUAGE_CODES.get(tgt_lang)

    if not src_code or not tgt_code:
        raise HTTPException(status_code=400, detail=f"Unsupported language: {src_lang} -> {tgt_lang}")

    try:
        # Load translation model
        translator = load_translator(src_code, tgt_code)

        # If indirect translation via English
        if isinstance(translator, tuple):
            translator1, translator2 = translator
            intermediate_text = "\n".join([translator1(chunk)[0]['translation_text'] for chunk in chunk_text(text)])
            translated_text = "\n".join([translator2(chunk)[0]['translation_text'] for chunk in chunk_text(intermediate_text)])

        else:
            translated_text = "\n".join([translator(chunk)[0]['translation_text'] for chunk in chunk_text(text)])

        return {"translated_text": translated_text}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Internal error: {str(e)}")
